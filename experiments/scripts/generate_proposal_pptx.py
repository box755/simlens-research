"""
Generate SimLens Proposal PPTX following SAILY Lab template style.
Output: SimLens_Proposal_2026-04-27.pptx (in project root)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = Path(__file__).resolve().parents[2] / "SimLens_Proposal_2026-04-27.pptx"

LAB_NAME = "Security and Artificial Intelligence Laboratory  (SAILY) Lab,\nNational Central University, Taoyuan, Taiwan"
DATE_STR = "2026/4/27"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

DARK = RGBColor(0x1F, 0x1F, 0x1F)
RED = RGBColor(0xC0, 0x00, 0x00)
GREY = RGBColor(0x80, 0x80, 0x80)
LIGHT_GREY = RGBColor(0xC0, 0xC0, 0xC0)
ACCENT = RGBColor(0x8B, 0x1A, 0x1A)  # 暗紅，模仿範本紅線
BLUE_BG = RGBColor(0xD9, 0xE7, 0xF5)
HEADER_BLUE = RGBColor(0x4A, 0x7E, 0xB6)


def add_footer(slide, page_num):
    # 左下日期
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(2), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = DATE_STR
    p.runs[0].font.size = Pt(10)
    p.runs[0].font.color.rgb = GREY
    # 中下 lab name
    tb = slide.shapes.add_textbox(Inches(3.5), Inches(7.05), Inches(7), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = LAB_NAME
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for r in para.runs:
            r.font.size = Pt(9)
            r.font.color.rgb = GREY
    # 右下頁碼
    tb = slide.shapes.add_textbox(Inches(12.7), Inches(7.05), Inches(0.5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = str(page_num)
    p.runs[0].font.size = Pt(10)
    p.runs[0].font.color.rgb = GREY
    p.alignment = PP_ALIGN.RIGHT


def add_title_with_redline(slide, title_text):
    # 標題文字
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.7))
    tf = tb.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.runs[0].font.size = Pt(32)
    p.runs[0].font.bold = True
    p.runs[0].font.name = "Times New Roman"
    p.runs[0].font.color.rgb = DARK
    # 紅線
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.95), Inches(12.5), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_bullets(slide, bullets, left=Inches(0.6), top=Inches(1.3), width=Inches(12), height=Inches(5.4),
                base_size=18, indent_size=14):
    """bullets: list of (level, text, color_or_None)"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            if len(item) == 3:
                level, text, color = item
            else:
                level, text = item
                color = None
        else:
            level, text, color = 0, item, None
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = level
        p.text = ("• " if level == 0 else "– ") + text
        for r in p.runs:
            r.font.name = "Times New Roman"
            r.font.size = Pt(base_size if level == 0 else indent_size)
            if color is not None:
                r.font.color.rgb = color
            else:
                r.font.color.rgb = DARK if level == 0 else RGBColor(0x40, 0x40, 0x40)
        p.space_after = Pt(6)


def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


# ───────────────────────── build ─────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Slide 1: Title
s = slide_blank(prs)
tb = s.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.3), Inches(2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Distillation + RLAIH for Segment-Level\nPersona-Conditioned Video Audience Simulation"
p.alignment = PP_ALIGN.CENTER
for r in p.runs:
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.name = "Times New Roman"
    r.font.color.rgb = DARK

tb = s.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12.3), Inches(1.2))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "傅聖祐 (Sheng-You Fu)"
p.alignment = PP_ALIGN.CENTER
p.runs[0].font.size = Pt(24); p.runs[0].font.name = "Times New Roman"; p.runs[0].font.color.rgb = DARK
p2 = tf.add_paragraph()
p2.text = "Advisor: Prof. Chia-Yu Lin"
p2.alignment = PP_ALIGN.CENTER
p2.runs[0].font.size = Pt(22); p2.runs[0].font.name = "Times New Roman"; p2.runs[0].font.color.rgb = DARK
add_footer(s, 1)

# Slide 2: Outline
s = slide_blank(prs)
add_title_with_redline(s, "Outline")
add_bullets(s, [
    (0, "Introduction", None),
    (1, "Motivation", LIGHT_GREY),
    (1, "Problem Description", LIGHT_GREY),
    (1, "Goal", LIGHT_GREY),
    (0, "Related Works", None),
    (0, "Proposed Framework", None),
    (0, "Experiments", None),
    (0, "Conclusion", None),
], top=Inches(1.4), base_size=22, indent_size=18)
add_footer(s, 2)

# Slide 3: Motivation
s = slide_blank(prs)
add_title_with_redline(s, "Motivation")
add_bullets(s, [
    (0, "Short-form video creators rely heavily on audience feedback to optimize content."),
    (1, "Traditional user studies are costly and slow."),
    (1, "Platform analytics (e.g., YouTube Analytics) only provide post-publish data, with days-to-weeks delay."),
    (0, "Existing AI-driven tools generate one comment per whole video."),
    (1, "They cannot reveal which segments engage which viewers."),
    (1, "Creators need actionable, time-localized insights — not aggregate sentiment."),
    (0, "Goal: shift audience research from passive whole-video feedback to active segment-level prediction."),
], top=Inches(1.3), base_size=18, indent_size=15)
add_footer(s, 3)

# Slide 4: Problem Description
s = slide_blank(prs)
add_title_with_redline(s, "Problem Description")
add_bullets(s, [
    (0, "Can a 3B small model — without real persona-level viewing data —"),
    (1, "generate persona-specific reactions for each segment of a 2-minute video?"),
    (0, "Three core challenges:"),
    (1, "C1. Data scarcity: no ground-truth dataset of (segment, persona, comment) tuples exists."),
    (1, "C2. Temporal granularity: model must reason about WHICH segment, not the whole video."),
    (1, "C3. On-device constraint: deployment target is a single 24GB consumer GPU."),
    (0, "Existing solutions (SimTube, Claude zero-shot, GPT-4o) cannot satisfy all three simultaneously."),
], top=Inches(1.3), base_size=18, indent_size=15)
add_footer(s, 4)

# Slide 5: Goal
s = slide_blank(prs)
add_title_with_redline(s, "Goal")
add_bullets(s, [
    (0, "Build SimLens — a segment-level persona-conditioned video audience simulator."),
    (0, "Quantitative KPIs:"),
    (1, "12 segments × 8 personas = 96-cell reaction matrix per 2-minute video."),
    (1, "Latency ≤ 15s on RTX 3090 (vs. ~120s for Claude API)."),
    (1, "Persona Consistency ≥ 0.80 (target: surpass Claude teacher 0.74)."),
    (1, "Segment Alignment Accuracy ≥ 80% (SimTube baseline: not applicable)."),
    (1, "API cost: $0 at inference time (vs. ~$94 / 1K calls for SimTube)."),
    (0, "Three contributions:"),
    (1, "C1 (System): first segment-level persona-conditioned video simulator."),
    (1, "C2 (Method): two-stage Distillation + RLAIH with 6-aspect reward."),
    (1, "C3 (Empirical): 3B student matches/exceeds 600B-class teacher on domain metrics."),
], top=Inches(1.2), base_size=17, indent_size=14)
add_footer(s, 5)

# Slide 6: Related Work — SimTube
s = slide_blank(prs)
add_title_with_redline(s, "Related Work — SimTube")
add_bullets(s, [
    (0, "Title: SimTube: Generating Simulated Video Comments through Multimodal AI", None),
    (0, "Problem: predict YouTube viewer comments before publishing.", None),
    (0, "Approach: Whisper + GPT-4 (visual captions) → Claude generates one comment per persona for the whole video.", None),
    (0, "Limitation: only whole-video output — no temporal granularity, no segment-level analysis.", RED),
    (1, "Cannot tell creators WHICH part of the video engages which audience.", RED),
    (1, "Relies on heavyweight closed-source APIs (Claude + GPT-4) → high cost, no on-device option.", RED),
], top=Inches(1.3), base_size=17, indent_size=14)
# reference footnote
tb = s.shapes.add_textbox(Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.4))
p = tb.text_frame.paragraphs[0]
p.text = "[1] Hung et al. SimTube: Generating Simulated Video Comments through Multimodal AI. arXiv:2411.09577, 2024."
p.runs[0].font.size = Pt(10); p.runs[0].font.name = "Times New Roman"; p.runs[0].font.color.rgb = GREY
add_footer(s, 6)

# Slide 7: Related Work — UMaT
s = slide_blank(prs)
add_title_with_redline(s, "Related Work — UMaT")
add_bullets(s, [
    (0, "Title: Everything Can Be Described in Words: A Unified Multi-Modal Framework with Semantic and Temporal Alignment", None),
    (0, "Problem: long-video understanding under multimodal heterogeneity and limited LLM context.", None),
    (0, "Approach: convert visual + audio into structured timestamped text segments → unified text representation for LLM.", None),
    (0, "Limitation: focuses on retrieval / QA — does NOT generate persona-conditioned reactions.", RED),
    (1, "Provides temporal alignment recipe but no audience modeling.", RED),
    (0, "SimLens borrows UMaT's structured-text alignment, extends it to generation + persona conditioning.", None),
], top=Inches(1.3), base_size=17, indent_size=14)
tb = s.shapes.add_textbox(Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.4))
p = tb.text_frame.paragraphs[0]
p.text = "[2] Bi & Xu. UMaT: Unified Multi-Modal Framework with Semantic and Temporal Alignment. arXiv:2503.09081, 2025."
p.runs[0].font.size = Pt(10); p.runs[0].font.name = "Times New Roman"; p.runs[0].font.color.rgb = GREY
add_footer(s, 7)

# Slide 8: Work Comparison
s = slide_blank(prs)
add_title_with_redline(s, "Work Comparison")
sub = s.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(12), Inches(0.45))
p = sub.text_frame.paragraphs[0]
p.text = "Comparing SimLens with state-of-the-art systems across four key dimensions."
p.runs[0].font.size = Pt(14); p.runs[0].font.italic = True; p.runs[0].font.color.rgb = RED; p.runs[0].font.name = "Times New Roman"

rows = [
    ["Method", "Segment-level reactions", "Persona-conditioned", "On-device (≤ 24GB)", "None-reaction modeling"],
    ["SimTube [1]", "", "✓", "", ""],
    ["Claude zero-shot",  "", "✓", "", ""],
    ["GPT-4o zero-shot", "", "✓", "", ""],
    ["Llama-3.2-3B zero-shot", "", "", "✓", ""],
    ["SimLens (ours)", "✓", "✓", "✓", "✓"],
]
n_rows, n_cols = len(rows), len(rows[0])
table_left, table_top = Inches(0.6), Inches(1.7)
table_w, table_h = Inches(12.1), Inches(4.5)
table_shape = s.shapes.add_table(n_rows, n_cols, table_left, table_top, table_w, table_h)
tbl = table_shape.table
col_widths = [Inches(2.6), Inches(2.6), Inches(2.3), Inches(2.3), Inches(2.3)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

for r_idx, row in enumerate(rows):
    for c_idx, val in enumerate(row):
        cell = tbl.cell(r_idx, c_idx)
        cell.text = val
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            for run in para.runs:
                run.font.name = "Times New Roman"
                run.font.size = Pt(15)
                if r_idx == 0:
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                elif r_idx == n_rows - 1:
                    run.font.bold = True
        if r_idx == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = HEADER_BLUE
        elif r_idx == n_rows - 1:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xCC)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = BLUE_BG if r_idx % 2 == 0 else RGBColor(0xEC, 0xF2, 0xF8)

# refs
tb = s.shapes.add_textbox(Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.4))
p = tb.text_frame.paragraphs[0]
p.text = "[1] Hung et al. SimTube. arXiv:2411.09577, 2024."
p.runs[0].font.size = Pt(10); p.runs[0].font.name = "Times New Roman"; p.runs[0].font.color.rgb = GREY
add_footer(s, 8)

# Slide 9: Proposed Framework section title
s = slide_blank(prs)
tb = s.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.3), Inches(2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "SimLens: Distillation + RLAIH for\nSegment-Level Persona Reaction Generation"
p.alignment = PP_ALIGN.CENTER
for r in p.runs:
    r.font.size = Pt(32); r.font.bold = True; r.font.name = "Times New Roman"; r.font.color.rgb = DARK
add_footer(s, 9)

# Slide 10: Architecture diagram placeholder
s = slide_blank(prs)
add_title_with_redline(s, "The Architecture of SimLens")
# placeholder box
ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.4), Inches(11.3), Inches(4.8))
ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
ph.line.color.rgb = LIGHT_GREY
tf = ph.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "[ Insert architecture diagram here ]"
p.alignment = PP_ALIGN.CENTER
p.runs[0].font.size = Pt(20); p.runs[0].font.italic = True; p.runs[0].font.color.rgb = GREY
p2 = tf.add_paragraph()
p2.text = "Replace with the modified architecture image\n(add: Llama-3B + LoRA, Claude as teacher, GPT-4o spot-check, Distillation, RLAIH)"
p2.alignment = PP_ALIGN.CENTER
for r in p2.runs:
    r.font.size = Pt(14); r.font.color.rgb = GREY; r.font.name = "Times New Roman"

# caption below
tb = s.shapes.add_textbox(Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.6))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Stage A: Whisper + LLaVA-NeXT (UMaT-style temporal alignment) → 12 enriched segments  |  "\
         "Stage B: Llama-3B + 8 persona-LoRA via Distillation (Claude) + RLAIH (Qwen judge, DPO)  |  "\
         "Stage C: same Llama-3B base → sentiment + report"
for r in p.runs:
    r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = DARK; r.font.name = "Times New Roman"
add_footer(s, 10)

# Slide 11: Method — Phase 1 Distillation
s = slide_blank(prs)
add_title_with_redline(s, "Method (1/2) — Phase 1: Knowledge Distillation")
add_bullets(s, [
    (0, "Goal: transfer Claude-3.5 Sonnet's segment-level persona reaction ability to Llama-3.2-3B."),
    (0, "Why distillation?  No real (segment, persona, comment) data exists; teacher LLM is the only feasible oracle."),
    (0, "Pipeline:"),
    (1, "Collect 100 × 2-min YouTube videos → run Stage A → 100 × 12 enriched segments."),
    (1, "For each (video, segment, persona) tuple: prompt Claude with cumulative narrative + current segment + persona YAML."),
    (1, "Output: a 10–50 word comment OR exactly \"None\" — explicitly modeling no-reaction behavior."),
    (1, "Total: 9,600 cells → ~4,800 comments + ~4,800 None labels (cost ≈ $115)."),
    (0, "SFT training: 8 independent LoRA adapters (rank=8) on 4-bit GPTQ Llama-3.2-3B-Instruct."),
    (0, "Novelty: training on \"None\" labels — teaches the model when a persona stays silent."),
], top=Inches(1.2), base_size=16, indent_size=14)
add_footer(s, 11)

# Slide 12: Method — Phase 2 RLAIH
s = slide_blank(prs)
add_title_with_redline(s, "Method (2/2) — Phase 2: RLAIH with 6-Aspect Reward")
add_bullets(s, [
    (0, "Goal: surpass Claude (teacher) on domain-specific dimensions via AI-feedback alignment."),
    (0, "Pipeline (per LoRA, iterated 2 rounds):"),
    (1, "Generate N=4 candidates per (segment, persona) prompt."),
    (1, "Score each candidate with 6 reward functions; build chosen/rejected preference pairs; run DPO."),
    (0, "6-aspect reward (with weights):"),
    (1, "R_relevance (25%): BERTScore + ROUGE-1 vs. current segment."),
    (1, "R_persona_consistency (20%) — Qwen3-32B-Q4 judge, PersonaGym rubric."),
    (1, "R_linguistic_habits (20%) — Qwen judge."),
    (1, "R_segment_relevance (15%) — NEW: penalizes comments matching wrong segments."),
    (1, "R_coherence (10%) — Qwen judge."),
    (1, "R_engagingness (10%) — UniEval."),
    (0, "Robustness: multi-judge ensemble (Qwen + Gemma2 + Llama-3.1) + GPT-4o spot-check on 200 samples."),
], top=Inches(1.2), base_size=15, indent_size=13)
add_footer(s, 12)

# Slide 13: Experiment — Data + Metrics
s = slide_blank(prs)
add_title_with_redline(s, "Experiments — Data & Evaluation Metrics")
add_bullets(s, [
    (0, "Data: 100 YouTube short-form videos (~2 min each), 5 categories × 20 videos."),
    (1, "Vlog / Tech Review / Food / Education / Entertainment; English; ≥ 10K views; official captions."),
    (1, "9,600 (video × segment × persona) cells; 85/10/5 train/val/test split."),
    (0, "Automatic metrics (Layer 2):"),
    (1, "BERTScore F1, ROUGE-1, Distinct-1/2 (NLG diversity)."),
    (1, "Persona Consistency, Linguistic Habits, Coherence, Engagingness (LLM-as-judge)."),
    (1, "Segment Alignment Accuracy, None Prediction F1 (NEW — SimLens-specific)."),
    (0, "Human evaluation (Layer 1):"),
    (1, "25 annotators × 8 videos × 3 personas; 7-point Likert on Relevance / Believability / Helpfulness."),
    (1, "Segment-localization quiz: given a generated comment, identify which of 12 segments it belongs to."),
    (0, "GPT-4o spot-check (200 samples) — verify local-judge ↔ GPT-4o Spearman ρ > 0.7."),
], top=Inches(1.2), base_size=15, indent_size=13)
add_footer(s, 13)

# Slide 14: Expected Results + Ablation
s = slide_blank(prs)
add_title_with_redline(s, "Expected Results & Ablation Study")

# table title
tb = s.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(12), Inches(0.4))
p = tb.text_frame.paragraphs[0]
p.text = "Expected main result (selected metrics) — values are projections from research plan."
p.runs[0].font.size = Pt(13); p.runs[0].font.italic = True; p.runs[0].font.color.rgb = RED; p.runs[0].font.name = "Times New Roman"

rows = [
    ["Method", "Persona Cons.", "Linguistic", "Segment Relev.", "None F1"],
    ["Llama-3B zero-shot", "0.42", "0.38", "0.40", "0.41"],
    ["Claude-3.5 Sonnet (teacher)", "0.74", "0.68", "0.65", "0.62"],
    ["GPT-4o zero-shot", "0.76", "0.70", "0.66", "—"],
    ["SimLens — Phase 1 only (SFT)", "0.71", "0.66", "0.62", "0.65"],
    ["SimLens — full (SFT + DPO) ★", "0.83", "0.81", "0.78", "0.78"],
]
n_rows, n_cols = len(rows), len(rows[0])
table_shape = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(1.5), Inches(8.5), Inches(3.0))
tbl = table_shape.table
widths = [Inches(2.7), Inches(1.5), Inches(1.4), Inches(1.5), Inches(1.4)]
for i, w in enumerate(widths):
    tbl.columns[i].width = w
for r_idx, row in enumerate(rows):
    for c_idx, val in enumerate(row):
        cell = tbl.cell(r_idx, c_idx)
        cell.text = val
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            for run in para.runs:
                run.font.name = "Times New Roman"; run.font.size = Pt(12)
                if r_idx == 0:
                    run.font.bold = True; run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                elif r_idx == n_rows - 1:
                    run.font.bold = True
        if r_idx == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = HEADER_BLUE
        elif r_idx == n_rows - 1:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xCC)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = BLUE_BG if r_idx % 2 == 0 else RGBColor(0xEC, 0xF2, 0xF8)

# Ablation list on the right
tb = s.shapes.add_textbox(Inches(9.3), Inches(1.5), Inches(3.8), Inches(5.0))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Planned ablations:"
p.runs[0].font.size = Pt(14); p.runs[0].font.bold = True; p.runs[0].font.name = "Times New Roman"; p.runs[0].font.color.rgb = DARK
items = [
    "A2 — w/o Phase 2 RLAIH",
    "A3 — w/o Phase 1 Distillation",
    "A4 — single LoRA (no per-persona)",
    "A5 — single-aspect reward",
    "A6 — w/o Segment Relevance reward",
    "A7 — w/o None handling",
    "A8 — w/o iterative DPO (1 round)",
]
for it in items:
    pp = tf.add_paragraph()
    pp.text = "• " + it
    pp.runs[0].font.size = Pt(11); pp.runs[0].font.name = "Times New Roman"; pp.runs[0].font.color.rgb = DARK
    pp.space_after = Pt(2)

# Bottom note
tb = s.shapes.add_textbox(Inches(0.6), Inches(4.7), Inches(12.1), Inches(1.8))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Why we expect SimLens to surpass Claude on domain metrics:"
p.runs[0].font.size = Pt(14); p.runs[0].font.bold = True; p.runs[0].font.name = "Times New Roman"; p.runs[0].font.color.rgb = RED
items = [
    "Per-persona LoRA adapters specialize beyond what a generalist teacher can express.",
    "RLAIH with domain-specific rewards directly optimizes Persona Consistency and Segment Relevance.",
    "Including None reactions teaches abstention behavior absent in zero-shot teachers.",
]
for it in items:
    pp = tf.add_paragraph()
    pp.text = "• " + it
    pp.runs[0].font.size = Pt(13); pp.runs[0].font.name = "Times New Roman"; pp.runs[0].font.color.rgb = DARK
    pp.space_after = Pt(2)
add_footer(s, 14)

# Slide 15: Conclusion
s = slide_blank(prs)
add_title_with_redline(s, "Conclusion")
add_bullets(s, [
    (0, "What we propose:"),
    (1, "SimLens — a 3B segment-level persona-conditioned video audience simulator,"),
    (1, "trained via two-stage Distillation + RLAIH with a novel 6-aspect reward."),
    (0, "What we expect to achieve:"),
    (1, "C1 (System): first end-to-end pipeline producing a 12 × 8 reaction matrix per 2-min video, runnable on 24GB GPU."),
    (1, "C2 (Method): demonstrate that distillation + RLAIH with a Segment Relevance reward can match or surpass a 600B-class teacher."),
    (1, "C3 (Empirical): None-reaction modeling — first system explicitly predicting persona silence (target F1 ≥ 0.78)."),
    (0, "Benefits (效益):"),
    (1, "Enables creators to optimize content BEFORE publishing, condensing weeks of feedback into minutes."),
    (1, "On-device, $0 inference cost — democratizes pre-publish audience analysis for small creators."),
    (1, "Foundation for future work: scaling to longer videos, more personas, real-time interaction."),
], top=Inches(1.2), base_size=16, indent_size=13)
add_footer(s, 15)

# Slide 16: References
s = slide_blank(prs)
add_title_with_redline(s, "References")
refs = [
    "[1] Hung et al. SimTube: Generating Simulated Video Comments through Multimodal AI. arXiv:2411.09577, 2024.",
    "[2] Bi & Xu. Everything Can Be Described in Words: A Unified Multi-Modal Framework with Semantic and Temporal Alignment (UMaT). arXiv:2503.09081, 2025.",
    "[3] Samuel et al. PersonaGym: Evaluating Persona Agents and LLMs. EMNLP 2025 Findings.",
    "[4] Rafailov et al. Direct Preference Optimization (DPO). NeurIPS 2023.",
    "[5] Lee et al. RLAIF: Scaling Reinforcement Learning from Human Feedback with AI Feedback. arXiv 2023, Google DeepMind.",
    "[6] Williams et al. MORLAIF: Multi-Objective Reinforcement Learning from AI Feedback. arXiv:2406.07496, 2024.",
    "[7] Yuan et al. Self-Rewarding Language Models. arXiv:2401.10020, 2024.",
    "[8] Hu et al. LoRA: Low-Rank Adaptation of Large Language Models. ICLR 2022.",
    "[9] Liu et al. LLaVA-NeXT: Improved Reasoning, OCR, and World Knowledge, 2024.",
    "[10] Radford et al. Robust Speech Recognition via Large-Scale Weak Supervision (Whisper). 2022.",
    "[11] Meta AI. Llama 3.2 Release Notes, 2024.",
    "[12] Lambert et al. Tülu 3: Pushing the Frontier of Open Language Model Post-Training, 2024.",
]
tb = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.6))
tf = tb.text_frame; tf.word_wrap = True
for i, r in enumerate(refs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = r
    for run in p.runs:
        run.font.size = Pt(13); run.font.name = "Times New Roman"; run.font.color.rgb = DARK
    p.space_after = Pt(4)
add_footer(s, 16)

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
